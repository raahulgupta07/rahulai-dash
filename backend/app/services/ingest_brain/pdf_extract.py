"""F09 EXTRACT (PDF / Word / PPT) — born-digital, GPU-free (Phase 2).

Text PDF  → pdfplumber tables + prose (camelot used opportunistically if present;
            it is NOT a hard dependency — missing camelot just means fewer rescued
            tables, never an error).
Word/PPT  → python-docx / python-pptx table iteration + prose.

Every reader is lazy-imported and fail-soft: a missing library or a parse error
returns ([], []) (or what it has) and is recorded as a note — never raises into
the ingest path. Vision (scanned/image) lives in ``vision_extract.py``.

Returns (list[RawTable], list[ProseBlock]).
"""
from __future__ import annotations

import logging
import re
from typing import Any, List, Tuple

from app.services.ingest_brain.contract import RawTable, ProseBlock

logger = logging.getLogger(__name__)


def _clean(v: Any) -> Any:
    if v is None:
        return None
    s = str(v).replace("\n", " ").strip()
    return s


def _to_table(name: str, grid: List[List[Any]], source_file: str, page: int) -> RawTable | None:
    grid = [[_clean(c) for c in row] for row in grid if any(_clean(c) for c in row)]
    if len(grid) < 2:
        return None
    width = max(len(r) for r in grid)
    grid = [r + [None] * (width - len(r)) for r in grid]
    header = [(_clean(h) or f"col_{i+1}") for i, h in enumerate(grid[0])]
    # de-dup header
    seen, hdr = {}, []
    for h in header:
        if h in seen:
            seen[h] += 1; hdr.append(f"{h}_{seen[h]}")
        else:
            seen[h] = 0; hdr.append(h)
    rows = grid[1:]
    return RawTable(name=name, header=hdr, rows=rows, source_file=source_file,
                    sheet=f"page_{page}", region_bbox=None,
                    notes=[f"extracted from PDF page {page}"])


def extract_pdf(path: str, *, filename: str = "") -> Tuple[List[RawTable], List[ProseBlock]]:
    """Born-digital PDF → tables + prose. NEVER raises."""
    tables: List[RawTable] = []
    prose: List[ProseBlock] = []
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for pi, page in enumerate(pdf.pages, 1):
                try:
                    for ti, grid in enumerate(page.extract_tables() or []):
                        tbl = _to_table(f"pdf_p{pi}_t{ti+1}", grid, filename, pi)
                        if tbl:
                            tables.append(tbl)
                    txt = (page.extract_text() or "").strip()
                    if txt:
                        prose.append(ProseBlock(title=f"{filename} — page {pi}",
                                                body=txt, source_file=filename, page=pi))
                except Exception:  # noqa: BLE001 — skip a bad page, keep going
                    logger.exception("pdf page %s failed", pi)
                    continue
    except Exception:  # noqa: BLE001
        logger.exception("extract_pdf failed for %s", path)
    # opportunistic camelot rescue (optional dep)
    try:
        import camelot  # noqa: F401
        if not tables:
            cam = camelot.read_pdf(path, pages="1-end")
            for ci, t in enumerate(cam):
                grid = t.df.values.tolist()
                tbl = _to_table(f"pdf_camelot_{ci+1}", grid, filename, ci + 1)
                if tbl:
                    tbl.notes.append("camelot lattice/stream")
                    tables.append(tbl)
    except Exception:  # noqa: BLE001 — camelot absent or failed → fine
        pass
    return tables, prose


def extract_docx(path: str, *, filename: str = "") -> Tuple[List[RawTable], List[ProseBlock]]:
    """Word → tables + prose. NEVER raises."""
    tables: List[RawTable] = []
    prose: List[ProseBlock] = []
    try:
        import docx
        d = docx.Document(path)
        for ti, t in enumerate(d.tables):
            grid = [[_clean(c.text) for c in row.cells] for row in t.rows]
            tbl = _to_table(f"docx_t{ti+1}", grid, filename, 0)
            if tbl:
                tbl.sheet = None
                tbl.notes = [f"extracted from Word table {ti+1}"]
                tables.append(tbl)
        body = "\n".join(p.text for p in d.paragraphs if p.text and p.text.strip())
        if body.strip():
            prose.append(ProseBlock(title=filename, body=body.strip(), source_file=filename))
    except Exception:  # noqa: BLE001
        logger.exception("extract_docx failed for %s", path)
    return tables, prose


def extract_pptx(path: str, *, filename: str = "") -> Tuple[List[RawTable], List[ProseBlock]]:
    """PowerPoint → tables + slide-text prose. NEVER raises."""
    tables: List[RawTable] = []
    prose: List[ProseBlock] = []
    try:
        from pptx import Presentation
        prs = Presentation(path)
        for si, slide in enumerate(prs.slides, 1):
            texts: List[str] = []
            for shape in slide.shapes:
                if shape.has_table:
                    tbl_obj = shape.table
                    grid = [[_clean(c.text) for c in row.cells] for row in tbl_obj.rows]
                    tbl = _to_table(f"pptx_s{si}", grid, filename, si)
                    if tbl:
                        tbl.sheet = f"slide_{si}"
                        tbl.notes = [f"extracted from slide {si}"]
                        tables.append(tbl)
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
            if texts:
                prose.append(ProseBlock(title=f"{filename} — slide {si}",
                                        body="\n".join(texts), source_file=filename, page=si))
    except Exception:  # noqa: BLE001
        logger.exception("extract_pptx failed for %s", path)
    return tables, prose
