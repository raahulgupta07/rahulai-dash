"""F09 Ingest Brain, Phase F: extract images embedded in .xlsx/.pptx files.

Each embedded image (chart / figure / picture) is sent to a vision callable
which reads it into a short, plain-text data description. Descriptions are
wrapped as :class:`ProseBlock` so they flow into the same prose lane as the
rest of the ingest brain. Everything here is fully fail-soft: any error per
image is skipped and total failure returns empty results. NEVER raises.
"""

from __future__ import annotations

import logging
from typing import Callable, Optional

from app.services.ingest_brain.contract import ProseBlock

logger = logging.getLogger(__name__)

_CHART_PROMPT = (
    "This image is a chart or figure from a document. In <=60 words describe "
    "what it shows and list any concrete numbers/labels/series visible. If it "
    "is not a chart, say 'not a data chart'."
)


def extract_embedded_charts(
    path: str,
    ext: str,
    *,
    filename: str = "",
    vision_infer: Optional[Callable[[bytes, str], str]] = None,
    max_images: int = 8,
) -> tuple[list, list[ProseBlock]]:
    """Extract embedded images from an xlsx/pptx file and describe them.

    Returns ``(tables, prose)`` where ``tables`` is always an empty list of
    RawTable for now, and ``prose`` is a list of :class:`ProseBlock`, one per
    usable embedded image. ``vision_infer`` is a callable built elsewhere with
    signature ``vision_infer(image_bytes, prompt) -> str``; if ``None`` no work
    is done. Never raises.
    """
    if vision_infer is None:
        return ([], [])

    try:
        ext_norm = (ext or "").lower().lstrip(".")
        if ext_norm in ("xlsx", "xlsm"):
            images = _extract_xlsx_images(path, max_images)
        elif ext_norm == "pptx":
            images = _extract_pptx_images(path, max_images)
        else:
            return ([], [])
    except Exception:  # pragma: no cover - total failure is fail-soft
        logger.debug("chart_extract: total failure for %s", path, exc_info=True)
        return ([], [])

    prose: list[ProseBlock] = []
    for idx, img_bytes in enumerate(images, start=1):
        try:
            if not img_bytes:
                continue
            text = vision_infer(img_bytes, _CHART_PROMPT)
            if not text or not text.strip():
                continue
            if "not a data chart" in text.lower():
                continue
            prose.append(
                ProseBlock(
                    title=f"{filename} — embedded image {idx}",
                    body=text.strip(),
                    source_file=filename,
                )
            )
        except Exception:
            logger.debug(
                "chart_extract: skipping image %d in %s", idx, path, exc_info=True
            )
            continue

    return ([], prose)


def _extract_xlsx_images(path: str, max_images: int) -> list[bytes]:
    """Pull image bytes from worksheet ``_images`` (private openpyxl API)."""
    import openpyxl  # lazy import

    out: list[bytes] = []
    try:
        wb = openpyxl.load_workbook(path)
    except Exception:
        logger.debug("chart_extract: cannot open xlsx %s", path, exc_info=True)
        return out

    try:
        for ws in wb.worksheets:
            for img in getattr(ws, "_images", []) or []:
                if len(out) >= max_images:
                    return out
                try:
                    data = img._data()  # private API, guard with try/except
                    if data:
                        out.append(data)
                except Exception:
                    logger.debug("chart_extract: bad xlsx image", exc_info=True)
                    continue
    finally:
        try:
            wb.close()
        except Exception:
            pass
    return out


def _extract_pptx_images(path: str, max_images: int) -> list[bytes]:
    """Pull image bytes from picture shapes in each slide."""
    from pptx import Presentation  # lazy import

    out: list[bytes] = []
    try:
        prs = Presentation(path)
    except Exception:
        logger.debug("chart_extract: cannot open pptx %s", path, exc_info=True)
        return out

    for slide in prs.slides:
        for shape in slide.shapes:
            if len(out) >= max_images:
                return out
            try:
                is_picture = getattr(shape, "shape_type", None) == 13
                image = getattr(shape, "image", None)
                if not is_picture and image is None:
                    continue
                if image is None:
                    image = shape.image
                blob = image.blob
                if blob:
                    out.append(blob)
            except Exception:
                logger.debug("chart_extract: bad pptx shape", exc_info=True)
                continue
    return out
