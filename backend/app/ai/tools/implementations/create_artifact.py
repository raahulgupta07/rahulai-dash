import asyncio
import base64
import json
import logging
from pathlib import Path
from typing import AsyncIterator, Dict, Any, Type, List, Optional

import aiofiles
from pydantic import BaseModel
from sqlalchemy import select
from sqlalchemy.orm import selectinload

from app.ai.tools.base import Tool
from app.models.file import File
from app.models.report_file_association import report_file_association

logger = logging.getLogger(__name__)


from app.ai.tools.metadata import ToolMetadata
from app.ai.tools.schemas import (
    ToolEvent,
    ToolStartEvent,
    ToolProgressEvent,
    ToolEndEvent,
)
from app.ai.tools.schemas.create_artifact import CreateArtifactInput, CreateArtifactOutput
from app.ai.llm import LLM
from app.ai.llm.types import ImageInput, Message, TextDeltaEvent
from app.models.artifact import Artifact
from app.models.visualization import Visualization
from app.dependencies import async_session_maker
from app.services.thumbnail_service import ThumbnailService
from app.services.artifact_libs import get_inline_scripts
from app.ai.code_execution.pptx_executor import PptxCodeExecutor, PptxPreviewService
from sqlalchemy import desc
from app.ai.tools.implementations._sandbox_context import SANDBOX_RUNTIME_PROMPT
from app.ai.prompt_language import build_language_directive


class CreateArtifactTool(Tool):
    """Tool for generating React-based artifact code for dashboards.

    This tool generates standalone React/JSX code that renders visualizations
    using ECharts, styled with Tailwind CSS, and transpiled in-browser via Babel.

    The generated code runs in a sandboxed iframe and receives visualization
    data via window.ARTIFACT_DATA.
    """

    @property
    def metadata(self) -> ToolMetadata:
        return ToolMetadata(
            name="create_artifact",
            description=(
                "Create or fully rebuild artifacts (dashboards, pages, slide presentations) from visualizations. "
                "Use for: new dashboards, full redesigns, large layout changes, or when edit_artifact cannot handle the scope. "
                "Modes: 'page' for interactive dashboards with KPI cards, charts, and responsive grids; "
                "'slides' for presentation decks (exportable to PPTX). "
                "IMPORTANT: visualization_ids are required - find them in previous create_data tool results "
                "shown as 'viz_id: <uuid>' in the conversation history. "
                "Do NOT ask the user for URLs or IDs - extract them from the conversation context. "
                "Only visualizations with successful step status are included."
            ),
            category="action",
            version="1.0.0",
            input_schema=CreateArtifactInput.model_json_schema(),
            output_schema=CreateArtifactOutput.model_json_schema(),
            max_retries=1,
            timeout_seconds=120,
            idempotent=False,
            required_permissions=[],
            is_active=True,
            tags=["artifact",  "dashboard", "slides"],
            allowed_modes=["chat", "deep"],
        )

    @property
    def input_model(self) -> Type[BaseModel]:
        return CreateArtifactInput

    @property
    def output_model(self) -> Type[BaseModel]:
        return CreateArtifactOutput

    # Path to the sandbox HTML file (relative to project root)
    # __file__ -> implementations -> tools -> ai -> app -> backend -> project_root
    SANDBOX_HTML_PATH = Path(__file__).parent.parent.parent.parent.parent.parent / "frontend" / "public" / "artifact-sandbox.html"

    async def _take_preview_screenshot(
        self,
        html_content: str,
    ) -> tuple[Optional[str], list[str]]:
        """Take a quick screenshot for planner reflection and capture JS errors.

        Returns (base64-encoded PNG string or None, list of JS error messages).
        """
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            return None, []

        js_errors: list[str] = []

        try:
            import tempfile, os
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                page = await browser.new_page(viewport={"width": 1280, "height": 720})

                # Capture JS errors during render
                page.on("pageerror", lambda err: js_errors.append(str(err)))

                # Write HTML to a temp file and navigate via file:// URL.
                # This allows vendored scripts (e.g. Tailwind runtime) that use
                # document.write() to work correctly — document.write fails on
                # about:blank pages used by set_content().
                tmp = tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="w", encoding="utf-8")
                try:
                    tmp.write(html_content)
                    tmp.close()
                    await page.goto(f"file://{tmp.name}", wait_until="networkidle")

                    # Wait for React to mount and charts to render (short timeout)
                    try:
                        await page.wait_for_function(
                            "window.__ARTIFACT_RENDER_COMPLETE__ === true",
                            timeout=8000,
                        )
                    except Exception:
                        pass  # Take screenshot anyway — partial render is still useful

                    await asyncio.sleep(0.3)
                    screenshot_bytes = await page.screenshot(type="png", full_page=False)
                    await browser.close()
                    return base64.b64encode(screenshot_bytes).decode("utf-8"), js_errors
                finally:
                    os.unlink(tmp.name)
        except Exception as e:
            logger.warning(f"Preview screenshot failed: {e}")
            return None, js_errors

    async def _generate_thumbnail_background(
        self,
        artifact_id: str,
        html_content: str,
        mode: str = "page",
    ) -> None:
        """Generate thumbnail in background and update artifact.

        Runs independently with its own database session.
        """
        try:
            thumbnail_service = ThumbnailService()
            thumbnail_path = await thumbnail_service.generate_thumbnail(
                artifact_id=artifact_id,
                html_content=html_content,
                mode=mode,
            )
            if thumbnail_path:
                # Use a fresh database session for the background update
                async with async_session_maker() as db:
                    from sqlalchemy import update
                    from app.models.artifact import Artifact
                    stmt = update(Artifact).where(Artifact.id == artifact_id).values(thumbnail_path=thumbnail_path)
                    await db.execute(stmt)
                    await db.commit()
        except Exception as e:
            logger.warning(f"Failed to generate thumbnail for artifact {artifact_id}: {e}")

    async def _load_completion_images(
        self,
        db: Any,
        head_completion_id: Optional[str],
    ) -> List[ImageInput]:
        """Load images attached to the head completion as ImageInput objects.

        Args:
            db: Database session
            head_completion_id: The completion ID to load images for

        Returns:
            List of ImageInput objects ready for vision-capable LLM
        """
        if not head_completion_id:
            return []

        images: List[ImageInput] = []
        try:
            # Query files associated with this completion that are images
            result = await db.execute(
                select(File)
                .join(report_file_association, report_file_association.c.file_id == File.id)
                .where(report_file_association.c.completion_id == head_completion_id)
                .where(File.content_type.startswith("image/"))
            )
            image_files = result.scalars().all()

            for f in image_files:
                if not f.path:
                    continue
                try:
                    async with aiofiles.open(f.path, 'rb') as file:
                        content = await file.read()
                    images.append(ImageInput(
                        data=base64.b64encode(content).decode('utf-8'),
                        media_type=f.content_type or 'image/png',
                        source_type='base64'
                    ))
                except Exception as e:
                    logger.warning(f"Failed to load image file {f.id}: {e}")

        except Exception as e:
            logger.warning(f"Failed to query completion images: {e}")

        return images

    def _build_thumbnail_html(self, artifact_data: dict, code: str, mode: str = "page") -> str:
        """Build HTML for thumbnail generation in headless browser.

        Args:
            artifact_data: The data to inject as window.ARTIFACT_DATA
            code: The LLM-generated artifact code
            mode: 'page' for React dashboards, 'slides' for pure HTML presentations

        Returns:
            Complete HTML string ready for headless browser rendering
        """
        data_json = json.dumps(artifact_data, default=str)

        # Slides mode: pure HTML + Tailwind (no React/Babel)
        if mode == "slides":
            slides_scripts = get_inline_scripts(mode="slides")
            slides_template = """<!DOCTYPE html>
<html>
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  __SLIDES_SCRIPTS__
  <style>
    html, body { height: 100%; margin: 0; padding: 0; }
    body { font-family: system-ui, -apple-system, sans-serif; }
    .slide { transition: opacity 0.3s ease-in-out; }
  </style>
</head>
<body class="bg-slate-900">
  <script>
    window.ARTIFACT_DATA = __ARTIFACT_DATA_JSON__;
    window.__ARTIFACT_RENDER_COMPLETE__ = false;
    setTimeout(function() {
      window.__ARTIFACT_RENDER_COMPLETE__ = true;
    }, 500);
  </script>

  __LLM_GENERATED_CODE__
</body>
</html>"""
            return slides_template.replace("__SLIDES_SCRIPTS__", slides_scripts).replace("__ARTIFACT_DATA_JSON__", data_json).replace("__LLM_GENERATED_CODE__", code)

        # Page mode: Build self-contained HTML mirroring ArtifactFrame.vue's approach.
        # get_inline_scripts("page") already includes all vendored libs + artifact-globals.js
        # so we only need to inject ARTIFACT_DATA, the LLM code, and render-complete detection.
        page_scripts = get_inline_scripts(mode="page")
        SC = '</' + 'script>'  # Avoid parser issues in this Python string too

        html = f"""<!DOCTYPE html>
<html>
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  {page_scripts}
  <style>
    html, body, #root {{ height: 100%; margin: 0; padding: 0; }}
    body {{ font-family: system-ui, -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; }}
  </style>
</head>
<body>
  <div id="root"></div>

  <script>
    window.ARTIFACT_DATA = {data_json};
    window.__ARTIFACT_RENDER_COMPLETE__ = false;
    window.__BOW_INFO = false;
  {SC}

  {code}

  <script>
    (function detectRenderComplete() {{
      var startTime = Date.now();
      var MAX_WAIT = 15000;
      function check() {{
        if (Date.now() - startTime > MAX_WAIT) {{
          window.__ARTIFACT_RENDER_COMPLETE__ = true;
          return;
        }}
        var root = document.getElementById('root');
        if (!root || root.children.length === 0) {{
          setTimeout(check, 200);
          return;
        }}
        var hasCharts = root.querySelectorAll('canvas').length > 0 ||
                        root.querySelectorAll('[_echarts_instance_]').length > 0;
        setTimeout(function() {{
          window.resizeAllCharts && window.resizeAllCharts();
          window.__ARTIFACT_RENDER_COMPLETE__ = true;
        }}, hasCharts ? 1500 : 300);
      }}
      setTimeout(check, 200);
    }})();
  {SC}
</body>
</html>"""
        return html

    async def _fix_code(
        self,
        code: str,
        errors: List[str],
        mode: str,
        runtime_ctx: Dict[str, Any],
        prompt_context: Dict[str, Any],
        screenshot_base64: Optional[str] = None,
        completion_images: Optional[List[ImageInput]] = None,
    ) -> str:
        """Attempt to fix code errors using the same prompt with error context.

        Args:
            code: The broken code
            errors: List of error messages
            mode: 'page' or 'slides'
            runtime_ctx: Runtime context for LLM access
            prompt_context: Context needed to rebuild the original prompt
                (user_prompt, title, viz_profiles, instructions_context,
                 report_title, allow_llm_see_data, messages_context, image_count)
            screenshot_base64: Optional screenshot of the broken render for visual context
            completion_images: Optional list of images from the head completion

        Returns:
            Fixed code string
        """
        error_text = "\n".join(f"- {e}" for e in errors[:5])  # Limit to first 5 errors

        # Rebuild the original prompt with full context
        base_prompt = self._build_prompt(
            user_prompt=prompt_context["user_prompt"],
            title=prompt_context["title"],
            mode=mode,
            viz_profiles=prompt_context["viz_profiles"],
            instructions_context=prompt_context["instructions_context"],
            report_title=prompt_context["report_title"],
            allow_llm_see_data=prompt_context["allow_llm_see_data"],
            messages_context=prompt_context.get("messages_context", ""),
            image_count=prompt_context.get("image_count", 0),
            organization_settings=prompt_context.get("organization_settings"),
        )

        # Build screenshot context if available
        screenshot_context = ""
        if screenshot_base64:
            screenshot_context = "\n\nA screenshot of the current broken render is attached. Use it to understand visual issues like layout problems, missing elements, or rendering errors."

        # Append error context and the broken code
        fix_prompt = f"""{base_prompt}

═══════════════════════════════════════════════════════════════════════════════
Fix the following errors
═══════════════════════════════════════════════════════════════════════════════

The previous code attempt produced these runtime errors:

{error_text}{screenshot_context}

Previous code:
```
{code}
```

Fix the errors while keeping the same design and functionality. Output the corrected code:"""

        # Skip fix if sigkill
        sigkill_event = runtime_ctx.get("sigkill_event")
        if sigkill_event and sigkill_event.is_set():
            return code

        # Use the same model for fixes
        llm = LLM(runtime_ctx.get("model"), usage_session_maker=async_session_maker)

        # Build image inputs: completion images + screenshot (if available)
        images: List[ImageInput] = []
        model = runtime_ctx.get("model")
        if model and getattr(model, "supports_vision", False):
            # Add completion images first (user's reference images)
            if completion_images:
                images.extend(completion_images)
            # Add screenshot of broken render last
            if screenshot_base64:
                images.append(ImageInput(data=screenshot_base64, media_type="image/png", source_type="base64"))

        try:
            chunks: list[str] = []
            async for evt in llm.inference_stream_v2(
                messages=[Message(role="user", content=fix_prompt)],
                images=images if images else None,
                usage_scope="create_artifact_fix",
            ):
                if isinstance(evt, TextDeltaEvent):
                    chunks.append(evt.text)
            response = "".join(chunks)
            return self._extract_code(response, mode=mode)
        except Exception as e:
            logger.exception("Error fixing code")
            # Return original code if fix fails
            return code

    def _build_viz_profile(self, viz: Dict[str, Any], allow_llm_see_data: bool) -> Dict[str, Any]:
        """Build a privacy-aware profile of a visualization's data."""
        # Enrich columns with dtype/unique_count/min/max from column_info (always — not sensitive)
        column_info = viz.get("column_info") or {}
        raw_columns = viz.get("columns", [])
        enriched_columns = []
        for c in raw_columns:
            col = dict(c) if isinstance(c, dict) else {"field": c}
            field = col.get("field") or col.get("headerName") or col.get("name")
            if field and field in column_info:
                meta = column_info[field]
                col["dtype"] = meta.get("dtype")
                col["unique_count"] = meta.get("unique_count")
                if meta.get("min") is not None:
                    col["min"] = meta["min"]
                if meta.get("max") is not None:
                    col["max"] = meta["max"]
            enriched_columns.append(col)

        profile: Dict[str, Any] = {
            "id": viz.get("id"),
            "title": viz.get("title"),
            "chart_type": viz.get("data_model_type") or "table",
            "row_count": viz.get("row_count", 0),
            "columns": enriched_columns,
        }

        # Include data model hints
        data_model = viz.get("dataModel") or {}
        if data_model:
            series = data_model.get("series", [])
            if series:
                profile["series_config"] = series[:3]  # First 3 series configs
            if data_model.get("group_by"):
                profile["group_by"] = data_model.get("group_by")

        # Include view configuration hints
        view = viz.get("view") or {}
        if view:
            inner_view = view.get("view") or view
            profile["view_config"] = {
                "type": inner_view.get("type"),
                "x": inner_view.get("x"),
                "y": inner_view.get("y"),
                "category": inner_view.get("category"),
                "value": inner_view.get("value"),
            }
            # Surface aggregation (top-level) + per-series aggregations so the
            # artifact can honor granular-data handling rather than reading
            # the first row.
            if inner_view.get("aggregation"):
                profile["view_config"]["aggregation"] = inner_view.get("aggregation")
            series_styles = inner_view.get("seriesStyles") or []
            series_aggs = [
                {"key": s.get("key"), "aggregation": s.get("aggregation")}
                for s in series_styles
                if isinstance(s, dict) and s.get("aggregation")
            ]
            if series_aggs:
                profile["view_config"]["series_aggregations"] = series_aggs
            default_filters = inner_view.get("defaultFilters") or []
            if default_filters:
                profile["view_config"]["default_filters"] = default_filters
            # Include palette if present
            palette = inner_view.get("palette") or {}
            if palette.get("colors"):
                profile["colors"] = palette.get("colors")[:5]

        # Include sample data if allowed
        if allow_llm_see_data:
            rows = viz.get("rows", [])
            if rows:
                profile["sample_rows"] = rows[:5]  # First 5 rows
                # Compute basic stats for numeric columns
                if rows and isinstance(rows[0], dict):
                    stats = {}
                    for col in viz.get("columns", []):
                        col_name = col if isinstance(col, str) else col.get("field", col.get("name"))
                        if col_name:
                            values = [r.get(col_name) for r in rows if r.get(col_name) is not None]
                            numeric_values = [v for v in values if isinstance(v, (int, float))]
                            if numeric_values:
                                stats[col_name] = {
                                    "min": min(numeric_values),
                                    "max": max(numeric_values),
                                    "sample_values": numeric_values[:3]
                                }
                            elif values:
                                unique = list(set(str(v) for v in values[:20]))
                                stats[col_name] = {
                                    "unique_count": len(unique),
                                    "sample_values": unique[:5]
                                }
                    if stats:
                        profile["column_stats"] = stats

        return profile

    async def run_stream(self, tool_input: Dict[str, Any], runtime_ctx: Dict[str, Any]) -> AsyncIterator[ToolEvent]:
        data = CreateArtifactInput(**tool_input)

        # Early validation: fail immediately if no visualization_ids provided
        if not data.visualization_ids or len(data.visualization_ids) == 0:
            yield ToolStartEvent(type="tool.start", payload={"title": data.title or "Artifact"})
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "error": "No visualization_ids provided. At least one visualization is required to create an artifact.",
                    },
                    "observation": {
                        "summary": "Failed to create artifact: no visualization_ids provided",
                        "error": {
                            "type": "validation_error",
                            "message": "visualization_ids is required and must contain at least one visualization ID. Create visualizations using create_data first, then use their IDs here.",
                        },
                    },
                },
            )
            return

        yield ToolStartEvent(type="tool.start", payload={"title": data.title or "Artifact"})
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "init"})

        # Get runtime context
        sigkill_event = runtime_ctx.get("sigkill_event")
        report = runtime_ctx.get("report")
        user = runtime_ctx.get("user")
        organization = runtime_ctx.get("organization")
        db = runtime_ctx.get("db")
        context_hub = runtime_ctx.get("context_hub")
        organization_settings = runtime_ctx.get("settings")

        # Check privacy setting
        allow_llm_see_data = True
        if organization_settings:
            try:
                allow_llm_see_data = organization_settings.get_config("allow_llm_see_data").value
            except Exception:
                allow_llm_see_data = True

        instruction_context_builder = runtime_ctx.get("instruction_context_builder") or (
            getattr(context_hub, "instruction_builder", None) if context_hub else None
        )

        # Get conversation history context (similar to create_data.py)
        context_view = runtime_ctx.get("context_view")
        messages_context = ""
        try:
            _messages_section_obj = getattr(context_view.warm, "messages", None) if context_view else None
            messages_context = _messages_section_obj.render() if _messages_section_obj else ""
        except Exception as e:
            logger.warning(f"Failed to extract messages context: {e}")
            messages_context = ""

        # Load images attached to the head completion for vision-capable models
        head_completion = runtime_ctx.get("head_completion")
        head_completion_id = str(head_completion.id) if head_completion else None
        completion_images = await self._load_completion_images(db, head_completion_id)

        # Validate model supports vision if images are present
        model = runtime_ctx.get("model")
        if completion_images and not getattr(model, "supports_vision", False):
            logger.info(f"Model doesn't support vision, skipping {len(completion_images)} completion images")
            completion_images = []

        # Note: Previous artifacts are now available via observation context (from create_artifact/read_artifact)
        # No need to fetch from DB - the planner can call read_artifact if needed

        # Fetch visualizations by ID from database
        visualizations: List[Dict[str, Any]] = []
        warnings: List[str] = []
        included_viz_ids: List[str] = []

        # Fetch all visualizations in a single batched query
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "loading_visualizations"})
        from app.models.query import Query
        from app.models.step import Step
        report_id = str(report.id) if report else None
        try:
            # populate_existing=True forces SQLAlchemy to refresh objects from DB
            # rather than returning stale identity-map copies (e.g. query.steps or
            # query.default_step may have been loaded before the step was created/updated)
            result = await db.execute(
                select(Visualization)
                .options(
                    selectinload(Visualization.query).selectinload(Query.default_step),
                    selectinload(Visualization.query).selectinload(Query.steps),
                )
                .where(Visualization.id.in_(data.visualization_ids))
                .execution_options(populate_existing=True)
            )
            fetched_vizs = {str(v.id): v for v in result.scalars().all()}
        except Exception as e:
            logger.exception("Failed to batch-fetch visualizations")
            fetched_vizs = {}
            warnings.append(f"Error fetching visualizations: {str(e)}")

        # Process each requested viz in order, validating and building entries
        for viz_id in data.visualization_ids:
            viz = fetched_vizs.get(viz_id)
            if viz is None:
                warnings.append(f"Visualization {viz_id} not found")
                continue

            # Validate viz belongs to the report
            if report_id and str(viz.report_id) != report_id:
                warnings.append(f"Visualization {viz_id} does not belong to this report")
                continue

            # Get the step with data (prefer default_step, fallback to latest step)
            step = None
            if viz.query and viz.query.default_step:
                step = viz.query.default_step
            elif viz.query and viz.query.steps:
                step = viz.query.steps[-1] if viz.query.steps else None

            # Check if the associated step is successful
            step_status = step.status if step else None
            if step_status != "success":
                _has_query = viz.query is not None
                _has_default = viz.query.default_step is not None if _has_query else False
                _steps_len = len(viz.query.steps) if _has_query and viz.query.steps else 0
                _default_step_id = getattr(viz.query, 'default_step_id', None) if _has_query else None
                logger.warning(
                    f"Visualization {viz_id} skipped: step_status='{step_status}', "
                    f"has_query={_has_query}, default_step_id={_default_step_id}, "
                    f"has_default_step={_has_default}, steps_count={_steps_len}"
                )
                warnings.append(f"Visualization {viz_id} skipped: step status is '{step_status or 'unknown'}' (not success)")
                continue

            # Get data directly from step (like frontend does)
            step_data = step.data if step else {}
            rows = (step_data.get("rows") or [])[:100] if step_data else []
            raw_columns = step_data.get("columns") or [] if step_data else []
            data_model = step.data_model if step else {}
            step_info = step_data.get("info") or {} if step_data else {}
            column_info = step_info.get("column_info") or {}

            # Keep raw column objects (with field/headerName) — matches the prompt contract
            columns = raw_columns

            # Extract field names for internal use (filterable columns, logging)
            column_fields = []
            for c in raw_columns:
                if isinstance(c, str):
                    column_fields.append(c)
                elif isinstance(c, dict):
                    col_name = c.get("field") or c.get("colId") or c.get("headerName") or c.get("name")
                    if col_name:
                        column_fields.append(col_name)

            # Build visualization entry
            view_dict = viz.view or {}
            query_id = str(viz.query_id) if viz.query_id else None

            ventry = {
                "id": str(viz.id),
                "title": viz.title,
                "query_id": query_id,
                "view": self._trim_none(view_dict),
                "data_model_type": (view_dict.get("view") or {}).get("type") or view_dict.get("type"),
                "columns": columns,
                "column_info": column_info,
                "row_count": len(rows),
                "rows": rows,
                "dataModel": data_model or {},
            }

            # Debug logging
            logger.info(f"Visualization {viz.title}: {len(rows)} rows, {len(column_fields)} columns: {column_fields[:5] if column_fields else 'none'}")
            if rows:
                logger.info(f"  Sample row keys: {list(rows[0].keys())[:5] if isinstance(rows[0], dict) else 'not a dict'}")

            visualizations.append(ventry)
            included_viz_ids.append(str(viz.id))

        # Early failure: if no valid visualizations were resolved, fail like create_data does with tables
        if not visualizations:
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "error": "No valid visualizations found. All requested visualization_ids were either not found, don't belong to this report, or have non-success step status.",
                    },
                    "observation": {
                        "summary": "Failed to create artifact: no valid visualizations resolved",
                        "error": {
                            "type": "no_valid_visualizations",
                            "message": "None of the requested visualization_ids could be used. Ensure visualizations exist, belong to this report, and have successful step status.",
                            "requested_ids": data.visualization_ids,
                            "warnings": warnings,
                        },
                    },
                },
            )
            return

        # Build visualization profiles (privacy-aware)
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "building_profiles"})
        viz_profiles = [self._build_viz_profile(v, allow_llm_see_data) for v in visualizations]

        # Emit visualizations_resolved
        yield ToolProgressEvent(type="tool.progress", payload={
            "stage": "visualizations_resolved",
            "tool_name": "create_artifact",
            "visualizations": [
                {"id": v["id"], "title": v["title"], "type": v.get("data_model_type", "")}
                for v in visualizations
            ],
        })

        # Build instruction context
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "building_context"})
        instructions_context = ""
        try:
            if instruction_context_builder is not None:
                inst_section = await instruction_context_builder.build(categories=["dashboard", "visualization", "general"])
                instructions_context = inst_section.render() or ""
        except Exception:
            pass

        # Create artifact early with pending status so frontend can show it
        artifact = Artifact(
            report_id=str(report.id) if report else None,
            user_id=str(user.id) if user else None,
            organization_id=str(organization.id) if organization else None,
            title=data.title or "Untitled Artifact",
            mode=data.mode,
            content={},  # Empty content initially
            generation_prompt=data.prompt,
            version=1,
            status="pending",
        )
        db.add(artifact)
        await db.commit()
        await db.refresh(artifact)

        # Notify frontend that artifact is created (pending)
        yield ToolProgressEvent(
            type="tool.progress",
            payload={
                "stage": "artifact_created",
                "artifact_id": str(artifact.id),
                "status": "pending",
                "timing": False,
            }
        )

        # Build the prompt for generating React code
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "building_prompt"})

        # Store prompt context for potential fix iterations
        prompt_context = {
            "user_prompt": data.prompt,
            "title": data.title,
            "viz_profiles": viz_profiles,
            "instructions_context": instructions_context,
            "report_title": getattr(report, 'title', None) if report else None,
            "allow_llm_see_data": allow_llm_see_data,
            "messages_context": messages_context,
            "image_count": len(completion_images),
            "organization_settings": organization_settings,
        }

        prompt = self._build_prompt(
            user_prompt=data.prompt,
            title=data.title,
            mode=data.mode,
            viz_profiles=viz_profiles,
            instructions_context=instructions_context,
            report_title=prompt_context["report_title"],
            allow_llm_see_data=allow_llm_see_data,
            messages_context=messages_context,
            image_count=len(completion_images),
            organization_settings=organization_settings,
        )

        # Stream from LLM
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "llm_generating"})
        llm = LLM(runtime_ctx.get("model"), usage_session_maker=async_session_maker)
        buffer = ""
        slides_detected = 0  # Track number of slides detected during streaming

        async for evt in llm.inference_stream_v2(
            messages=[Message(role="user", content=prompt)],
            images=completion_images if completion_images else None,
            usage_scope="create_artifact",
            usage_scope_ref_id=str(report.id) if report else None,
        ):
            if sigkill_event and sigkill_event.is_set():
                break
            if isinstance(evt, TextDeltaEvent):
                buffer += evt.text

            # For slides mode, detect new slides as they're generated
            if data.mode == "slides":
                # Count slide sections in buffer
                current_slides = buffer.count('<section class="slide"')
                if current_slides > slides_detected:
                    # New slide detected
                    for i in range(slides_detected, current_slides):
                        yield ToolProgressEvent(
                            type="tool.progress",
                            payload={
                                "stage": "slide_generated",
                                "slide_index": i,
                                "total_slides": current_slides,
                                "timing": False,
                            }
                        )
                    slides_detected = current_slides

            # Stream partial updates
            if len(buffer) % 100 == 0:  # Throttle updates
                yield ToolProgressEvent(
                    type="tool.progress",
                    payload={"stage": "generating", "chars": len(buffer), "timing": False}
                )

        # Check sigkill after LLM generation
        if sigkill_event and sigkill_event.is_set():
            # Update artifact to stopped status
            artifact.status = "stopped"
            await db.commit()
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {"success": False, "artifact_id": str(artifact.id), "error": "Stopped by user"},
                    "observation": {"summary": "Artifact creation stopped by user", "artifact_id": str(artifact.id), "stopped": True},
                },
            )
            return

        # Extract the code from the response
        code = self._extract_code(buffer, mode=data.mode)

        # ═══════════════════════════════════════════════════════════════════════
        # Mode-specific processing: slides uses python-pptx, page skips to save
        # ═══════════════════════════════════════════════════════════════════════

        pptx_path: Optional[str] = None
        pptx_success: bool = True
        pptx_error: Optional[str] = None
        preview_images: List[str] = []

        if data.mode == "slides":
            # ═══════════════════════════════════════════════════════════════════
            # SLIDES MODE: Execute python-pptx code and generate previews
            # ═══════════════════════════════════════════════════════════════════
            yield ToolProgressEvent(
                type="tool.progress",
                payload={"stage": "executing_pptx_code"}
            )

            try:
                # Prepare data for execution
                report_data = {
                    "id": str(report.id) if report else None,
                    "title": getattr(report, "title", None) if report else None,
                    "theme": getattr(report, "theme", None) if report else None,
                }

                # Setup output path
                uploads_dir = Path(__file__).parent.parent.parent.parent.parent / "uploads" / "pptx"
                uploads_dir.mkdir(parents=True, exist_ok=True)
                output_path = uploads_dir / f"{artifact.id}.pptx"

                # Execute the python-pptx code
                executor = PptxCodeExecutor(logger=logger)
                result_path, output_log = executor.execute_pptx_code(
                    code=code,
                    visualizations=visualizations,
                    report=report_data,
                    output_path=output_path,
                )

                pptx_path = str(result_path)

                yield ToolProgressEvent(
                    type="tool.progress",
                    payload={"stage": "generating_previews"}
                )

                # Generate preview images
                preview_service = PptxPreviewService(logger=logger)
                preview_images = preview_service.generate_previews(
                    pptx_path=result_path,
                    artifact_id=str(artifact.id),
                )

            except Exception as e:
                logger.error(f"PPTX execution failed: {e}")
                pptx_success = False
                pptx_error = str(e)

        yield ToolProgressEvent(type="tool.progress", payload={"stage": "saving_artifact"})

        # Build content object
        content: Dict[str, Any] = {
            "code": code,
            "visualization_ids": included_viz_ids,
        }

        # Add slides-specific content
        if data.mode == "slides" and preview_images:
            content["preview_images"] = preview_images

        # Update the pending artifact with content and mark as completed
        artifact.content = content
        artifact.status = "completed" if (data.mode != "slides" or pptx_success) else "failed"

        # On slides failure, surface the error on the artifact so the UI shows an
        # error panel (with a Fix button) instead of dumping raw python as text.
        if data.mode == "slides" and not pptx_success and pptx_error:
            artifact.render_errors = [pptx_error]

        # Set pptx_path for slides mode
        if pptx_path:
            artifact.pptx_path = pptx_path

        await db.commit()
        await db.refresh(artifact)

        # Page mode: take preview screenshot for planner reflection + generate thumbnail
        screenshot_base64: Optional[str] = None
        render_errors: list[str] = []
        if data.mode == "page":
            artifact_data = {
                "report": {
                    "id": str(report.id) if report else None,
                    "title": getattr(report, "title", None) if report else None,
                    "theme": getattr(report, "theme", None) if report else None,
                },
                "visualizations": visualizations,
            }
            thumbnail_html = self._build_thumbnail_html(artifact_data, code, mode=data.mode)

            # Take preview screenshot (synchronous, ~3-5s) if model supports vision
            model = runtime_ctx.get("model")
            if allow_llm_see_data and model and getattr(model, "supports_vision", False):
                yield ToolProgressEvent(type="tool.progress", payload={"stage": "capturing_preview"})
                screenshot_base64, render_errors = await self._take_preview_screenshot(thumbnail_html)

            # Persist screenshot and render errors on artifact for later retrieval (read_artifact)
            if screenshot_base64 or render_errors:
                artifact.screenshot_base64 = screenshot_base64
                artifact.render_errors = render_errors or None
                await db.commit()

            # Generate thumbnail in background (for stored thumbnail, non-blocking)
            asyncio.create_task(
                self._generate_thumbnail_background(
                    artifact_id=str(artifact.id),
                    html_content=thumbnail_html,
                    mode=data.mode,
                )
            )
        elif preview_images:
            # For slides mode, use the first preview image as thumbnail
            first_preview = Path(__file__).parent.parent.parent.parent.parent / "uploads" / preview_images[0]
            if first_preview.exists():
                artifact.thumbnail_path = preview_images[0]
                await db.commit()

        output = CreateArtifactOutput(
            artifact_id=str(artifact.id),
            code=code,
            mode=data.mode,
            title=data.title,
            version=artifact.version,
        ).model_dump()

        # Add UI preview fields (similar to read_artifact)
        code_lines = code.count('\n') + 1 if code else 0
        output["artifact_preview"] = {
            "artifact_id": str(artifact.id),
            "title": data.title or "Untitled",
            "mode": data.mode,
            "version": artifact.version,
            "code_stats": {
                "chars": len(code),
                "lines": code_lines,
            },
            "visualization_ids": included_viz_ids,
            "visualization_count": len(visualizations),
        }
        # Code for collapsible toggle (collapsed by default in UI)
        output["code_preview"] = {
            "language": "jsx",
            "code": code,
            "collapsed_default": True,
        }

        # Build observation message
        summary_msg = f"Created artifact '{data.title or 'Untitled'}' with {len(code)} characters of code"
        if data.mode == "slides" and preview_images:
            summary_msg += f". Generated {len(preview_images)} slide preview images."
        elif render_errors:
            summary_msg += f". RENDER FAILED with {len(render_errors)} error(s): {render_errors[0]}"
            if len(render_errors) > 1:
                summary_msg += f" (and {len(render_errors) - 1} more)"
            summary_msg += ". The dashboard code has a bug — use edit_artifact to fix the specific error."
        elif screenshot_base64:
            summary_msg += ". Screenshot of the rendered dashboard is attached — review it for visual correctness."

        observation: Dict[str, Any] = {
            "summary": summary_msg,
            "artifact_id": str(artifact.id),
            "mode": data.mode,
            "visualization_count": len(visualizations),
            "visualization_ids": included_viz_ids,
        }
        if render_errors:
            observation["render_errors"] = render_errors

        # Add preview screenshot for planner reflection (page mode)
        if screenshot_base64:
            observation["images"] = [{
                "data": screenshot_base64,
                "media_type": "image/png",
                "source_type": "base64",
            }]

        # Add slides-specific info
        if data.mode == "slides":
            if preview_images:
                observation["preview_images"] = preview_images
                observation["slide_count"] = len(preview_images)
            if pptx_path:
                observation["pptx_path"] = pptx_path

        if warnings:
            observation["warnings"] = warnings

        yield ToolEndEvent(
            type="tool.end",
            payload={
                "output": output,
                "observation": observation,
            }
        )

    def _trim_none(self, obj: Any) -> Any:
        """Remove None values and empty collections from nested structures."""
        try:
            if isinstance(obj, dict):
                out = {}
                for k, v in obj.items():
                    tv = self._trim_none(v)
                    if tv is None:
                        continue
                    if isinstance(tv, (dict, list)) and len(tv) == 0:
                        continue
                    out[k] = tv
                return out
            if isinstance(obj, list):
                items = [self._trim_none(v) for v in obj]
                return [v for v in items if not (v is None or (isinstance(v, (dict, list)) and len(v) == 0))]
            return obj
        except Exception:
            return obj

    def _build_slides_prompt(
        self,
        user_prompt: str,
        title: str | None,
        viz_profiles: List[Dict[str, Any]],
        instructions_context: str,
        report_title: str | None,
        allow_llm_see_data: bool,
        messages_context: str = "",
        image_count: int = 0,
        organization_settings: Any = None,
    ) -> str:
        """Build the prompt for generating slides using python-pptx code."""
        viz_json = json.dumps(viz_profiles, indent=2, default=str)

        language_directive = build_language_directive(organization_settings)

        # Build attached images context
        images_context = ""
        if image_count > 0:
            images_context = f"\n**Attached Images:** {image_count} image(s) provided for visual reference. Use these to understand the design intent, branding, color schemes, or layout preferences the user wants to incorporate."

        return f"""Role: presentation author using python-pptx.{language_directive}
Generate python-pptx code to create a polished slide deck.

═══════════════════════════════════════════════════════════════════════════════
AVAILABLE IN NAMESPACE (already provided — do not import)
═══════════════════════════════════════════════════════════════════════════════

Python-pptx classes and functions:
- Presentation, Inches, Pt, Emu, RGBColor
- PP_ALIGN, MSO_ANCHOR, MSO_SHAPE
- XL_CHART_TYPE, XL_LEGEND_POSITION
- CategoryChartData, ChartData

Note: Inches, Pt, Emu are functions, not methods.
   Use: Inches(1), Pt(24), Emu(914400)
   Not: 1.inches, 24.pt, value.inches

Data variables:
- visualizations: List[Dict] — each has 'title', 'columns', 'rows'
- report: Dict with 'id', 'title', 'theme'

Output:
- _pptx_output_path: str — path to save the presentation to

═══════════════════════════════════════════════════════════════════════════════
YOUR VISUALIZATIONS
═══════════════════════════════════════════════════════════════════════════════

{viz_json}

{"(Full sample data included above)" if allow_llm_see_data else "(Data samples hidden for privacy - use column names and row_count)"}

═══════════════════════════════════════════════════════════════════════════════
TASK
═══════════════════════════════════════════════════════════════════════════════

**Report Title:** {report_title or title or 'Presentation'}
**User Request:** {user_prompt}
{images_context}
{f"**Organization Instructions:** {instructions_context}" if instructions_context else ""}

═══════════════════════════════════════════════════════════════════════════════
PYTHON-PPTX QUICK REFERENCE
═══════════════════════════════════════════════════════════════════════════════

**Setup (16:9 widescreen):**
```python
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
```

**Add blank slide with dark background:**
```python
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42)  # slate-900
```

**Add text box:**
```python
txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Title Text"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
```

**Add bar chart (use this pattern for charts):**
```python
chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('Revenue', (1.2, 1.5, 1.8, 2.1))

x, y, cx, cy = Inches(1), Inches(2), Inches(11), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
).chart

# Style the chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
plot = chart.plots[0]
plot.has_data_labels = True

# ⚠️ MANDATORY: native pptx charts default to BLACK fonts → INVISIBLE on dark slides.
# Recolor EVERY chart's fonts to your slide's text color. Call this on every chart.
style_chart_text(chart, TEXT_LIGHT)   # TEXT_LIGHT on dark slides; TEXT_DARK on light slides
```

**`style_chart_text` helper — paste once, call on every chart:**
```python
def style_chart_text(chart, color):
    # Axis tick labels (category + value)
    for axis in (getattr(chart, 'category_axis', None), getattr(chart, 'value_axis', None)):
        if axis is None:
            continue
        try:
            axis.tick_labels.font.color.rgb = color
            axis.tick_labels.font.size = Pt(11)
        except Exception:
            pass
    # Legend
    if chart.has_legend:
        try:
            chart.legend.font.color.rgb = color
            chart.legend.font.size = Pt(11)
        except Exception:
            pass
    # Data labels on every plot/series
    for plot in chart.plots:
        if getattr(plot, 'has_data_labels', False):
            try:
                plot.data_labels.font.color.rgb = color
                plot.data_labels.font.size = Pt(10)
            except Exception:
                pass
```

**Other chart types:**
- XL_CHART_TYPE.COLUMN_CLUSTERED - vertical bars
- XL_CHART_TYPE.LINE - line chart
- XL_CHART_TYPE.PIE - pie chart
- XL_CHART_TYPE.AREA - area chart

**Dark background (slate-900 = RGB(15, 23, 42)):**
```python
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42)
```

**Access visualization data:**
```python
viz = visualizations[0]
columns = viz['columns']  # e.g. ['AlbumTitle', 'Revenue', 'UnitsSold']
rows = viz['rows']        # list of dicts like {{'AlbumTitle': 'Greatest Hits', 'Revenue': 1500.0}}

# Get categories and values for a chart:
categories = [str(row[columns[0]]) for row in rows]  # First column as labels
values = [float(row[columns[1]]) if row[columns[1]] else 0 for row in rows]  # Second column as values

# IMPORTANT: columns[i] returns a string like 'Revenue', then use that to index into row
# row[columns[1]] is the same as row['Revenue'] if columns[1] == 'Revenue'
```

═══════════════════════════════════════════════════════════════════════════════
DESIGN PHILOSOPHY - CREATE BEAUTIFUL, PROFESSIONAL SLIDES
═══════════════════════════════════════════════════════════════════════════════

**THEME PER TOPIC — pick light OR dark to fit the data, then commit to it:**
Choose the background mode that suits the subject (e.g. finance/executive → dark navy; consumer/music/health → clean light). Pick ONE mode for the whole deck; do not mix light and dark slides.
- Define `TEXT_DARK = RGBColor(15, 23, 42)` AND `TEXT_LIGHT = RGBColor(255, 255, 255)` up front.
- On a DARK background → use TEXT_LIGHT for titles, body, KPI numbers, and `style_chart_text(chart, TEXT_LIGHT)`.
- On a LIGHT background → use TEXT_DARK and `style_chart_text(chart, TEXT_DARK)`.

**⚠️ CONTRAST IS NON-NEGOTIABLE (this is the #1 failure):**
- Native python-pptx charts render their axis labels, legend, and data labels in BLACK by default. On a dark slide they are INVISIBLE. You MUST call `style_chart_text(chart, <slide text color>)` on EVERY chart you add — no exceptions.
- Never place dark text on a dark fill or light text on a light fill. Every textbox, KPI number, label, axis tick and legend must contrast its background.
- Accent shapes/cards: if the card fill is dark, its text is TEXT_LIGHT; if the card fill is a light accent, its text is TEXT_DARK.

**COLOR STRATEGY - Be Topic-Specific:**
Choose colors that feel designed for THIS topic. If your colors would work for any presentation, you haven't made specific enough choices.

Structure: One DOMINANT color (60-70% visual weight), 1-2 supporting tones, one accent.

Example palettes (pick one that fits the topic):
- **Midnight Executive**: Navy (0,31,63), Steel (119,136,153), Gold accent (212,175,55)
- **Forest & Moss**: Deep green (34,87,76), Sage (138,154,91), Cream (245,245,220)
- **Coral Energy**: Coral (255,127,80), Teal (0,128,128), Sand (244,232,214)
- **Ocean Depths**: Deep blue (0,51,102), Aqua (0,180,180), Pearl (240,248,255)
- **Sunset Warm**: Burgundy (128,0,32), Orange (255,140,0), Cream (255,253,240)
- **Modern Minimal**: Charcoal (54,69,79), Light gray (220,220,220), Teal accent (0,150,136)

**Layout variety — vary between slides:**
Every slide should have visual elements — charts, shapes, or decorative elements. Avoid text-only slides.

Vary layouts between:
- Two-column (text left, chart right or vice versa)
- Full-width chart with title above
- KPI cards in a row (3-4 metric boxes)
- Chart with callout boxes for key insights
- Split layout with accent shape dividers

**Typography:**
- Titles: 36-44pt bold, interesting positioning (not always centered)
- Body text: 18-24pt, left-aligned (avoid center-aligning body text)
- KPI numbers: 48-72pt bold for impact
- Use font color contrast: white on dark, dark on light accents

**VISUAL ELEMENTS TO ADD:**
- Accent shapes: rectangles, rounded rectangles for backgrounds
- Divider lines or shapes between sections
- Colored boxes behind KPI numbers
- Subtle shape overlays for visual interest

**Common mistakes to avoid:**
- Using `value.inches` instead of `Inches(value)` — Inches/Pt/Emu are functions.
- Repeating the same layout across slides — vary it.
- Center-aligning body text — use left alignment.
- Using only blue without topic-specific reasoning.
- Text-only slides without visual elements.
- Accent lines directly under titles (hallmark of generic slides).
- Cramming too much data — limit charts to top 8-10 items.

**Technical requirements:**
1. Define `generate_slides(visualizations, report)` returning a Presentation.
2. Use 16:9 widescreen: Inches(13.333) x Inches(7.5).
3. Create real charts with slide.shapes.add_chart() + CategoryChartData.
4. Use visualization data from the visualizations list.
5. Margins: start shapes at Inches(0.75) to Inches(1) from edges.

**⚠️ DATA SAFETY — NEVER add a chart with empty categories (this hard-crashes the
whole deck with "chart data contains no categories"):**
- Before building any `CategoryChartData`, derive categories and FILTER OUT rows
  whose label is None/empty: `cats = [str(r[label]) for r in rows if r.get(label) not in (None, "")]`.
- If `cats` is empty (a KPI / single-value / non-categorical visualization), DO
  NOT call `add_chart`. Render that visualization as a KPI card or a small text
  table instead (big number + label). A deck of KPI cards is fine.
- Coerce values defensively: `vals = [float(r.get(measure) or 0) for r in rows]`.
  Guard every `add_series` so categories and values are non-empty and equal length.
- One bad visualization must not sink the deck — wrap each chart's build in a
  try/except and fall back to a text/KPI slide for that one viz.

═══════════════════════════════════════════════════════════════════════════════
OUTPUT FORMAT - Example with Design Principles Applied
═══════════════════════════════════════════════════════════════════════════════

```python
def generate_slides(visualizations, report):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette - choose colors that fit the topic
    PRIMARY = RGBColor(0, 51, 102)      # Deep blue
    SECONDARY = RGBColor(0, 128, 128)   # Teal
    ACCENT = RGBColor(255, 140, 0)      # Orange accent
    BG_DARK = RGBColor(15, 23, 42)      # Dark background
    TEXT_LIGHT = RGBColor(255, 255, 255)
    TEXT_DARK = RGBColor(15, 23, 42)    # for light slides / light accent cards
    TEXT_MUTED = RGBColor(148, 163, 184)
    TEXT_ON_BG = TEXT_LIGHT             # flip to TEXT_DARK if you choose a light BG

    def style_chart_text(chart, color):
        for axis in (getattr(chart, 'category_axis', None), getattr(chart, 'value_axis', None)):
            if axis is None:
                continue
            try:
                axis.tick_labels.font.color.rgb = color
                axis.tick_labels.font.size = Pt(11)
            except Exception:
                pass
        if chart.has_legend:
            try:
                chart.legend.font.color.rgb = color
                chart.legend.font.size = Pt(11)
            except Exception:
                pass
        for plot in chart.plots:
            if getattr(plot, 'has_data_labels', False):
                try:
                    plot.data_labels.font.color.rgb = color
                    plot.data_labels.font.size = Pt(10)
                except Exception:
                    pass

    def set_background(slide, color=BG_DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_accent_shape(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 1: Title with accent shape
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    # Accent shape behind title
    add_accent_shape(slide, Inches(0), Inches(2.5), Inches(5), Inches(2.5), PRIMARY)

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = report.get('title', 'Presentation')
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 2: KPI Cards Row (if we have numeric data)
    # ═══════════════════════════════════════════════════════════════
    if visualizations and visualizations[0].get('rows'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)

        viz = visualizations[0]
        rows = viz.get('rows', [])
        columns = viz.get('columns', [])

        # Create 3 KPI cards across the slide
        card_width = Inches(3.5)
        card_height = Inches(2.5)
        start_x = Inches(1)
        card_y = Inches(2.5)
        gap = Inches(0.5)

        for i, col in enumerate(columns[:3]):
            if i >= 3:
                break
            x = start_x + i * (card_width + gap)

            # Card background
            card = add_accent_shape(slide, x, card_y, card_width, card_height, PRIMARY)

            # Value (large number)
            val = rows[0].get(col, 0) if rows else 0
            val_box = slide.shapes.add_textbox(x + Inches(0.3), card_y + Inches(0.5), card_width - Inches(0.6), Inches(1.2))
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = "{{:,.0f}}".format(float(val)) if isinstance(val, (int, float)) else str(val)
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = TEXT_LIGHT

            # Label
            label_box = slide.shapes.add_textbox(x + Inches(0.3), card_y + Inches(1.7), card_width - Inches(0.6), Inches(0.6))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = col
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MUTED

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 3: Chart with title (different layout)
    # ═══════════════════════════════════════════════════════════════
    if visualizations:
        viz = visualizations[0]
        columns = viz.get('columns', [])
        rows = viz.get('rows', [])

        if len(columns) >= 2 and rows:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_background(slide)

            # Title on left side
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(5), Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = viz.get('title', 'Data Analysis')
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = TEXT_LIGHT

            # Extract data
            col_label = columns[0]
            col_value = columns[1]
            categories = [str(row.get(col_label, ''))[:20] for row in rows[:8]]
            values = [float(row.get(col_value, 0) or 0) for row in rows[:8]]

            # Chart (full width below title)
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series(col_value, tuple(values))

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5),
                chart_data
            ).chart
            chart.has_legend = False
            style_chart_text(chart, TEXT_ON_BG)  # MANDATORY — else axis labels are black/invisible

    return prs

# Execute and save
prs = generate_slides(visualizations, report)
prs.save(_pptx_output_path)
```

Create a beautiful, varied presentation following these design principles. Each slide should look DIFFERENT from the others. Use visual elements, accent shapes, and thoughtful color choices:"""

    def _build_page_prompt(
        self,
        user_prompt: str,
        title: str | None,
        viz_profiles: List[Dict[str, Any]],
        instructions_context: str,
        report_title: str | None,
        allow_llm_see_data: bool,
        messages_context: str = "",
        image_count: int = 0,
        organization_settings: Any = None,
    ) -> str:
        """Build the prompt for generating page/dashboard (React + ECharts)."""
        viz_json = json.dumps(viz_profiles, indent=2, default=str)

        language_directive = build_language_directive(organization_settings)

        # Build attached images context
        images_context = ""
        if image_count > 0:
            images_context = f"\n**Attached Images:** {image_count} image(s) provided for visual reference. Use these to understand the design intent, branding, color schemes, or layout preferences the user wants to incorporate."

        # Note: Previous artifact code is now available via observation context (from create_artifact/read_artifact)
        # The planner can call read_artifact if needed to load previous code into context

        return f"""Role: frontend developer and data visualization engineer.

═══════════════════════════════════════════════════════════════════════════════
Design request (primary specification — takes precedence when it conflicts with defaults)
═══════════════════════════════════════════════════════════════════════════════

**Report Title:** {report_title or title or 'Dashboard'}
**User Request:** {user_prompt}
{images_context}
{f"**Organization Instructions:**{chr(10)}{instructions_context}" if instructions_context else ""}

{f"**Conversation History:**{chr(10)}{messages_context}" if messages_context else ""}
{language_directive}

If the user specified a theme, layout, colors, or style above — follow that exactly.
If the user did not specify styling, use the design guidance at the end of this prompt.

═══════════════════════════════════════════════════════════════════════════════
REFERENCE — TOOLS, COMPONENTS & DATA
═══════════════════════════════════════════════════════════════════════════════

{SANDBOX_RUNTIME_PROMPT}

CHARTING:

**`<EChart height={{N}} option={{{{...}}}} />`** — chart wrapper. Supports ALL ECharts chart types. 'dash' theme pre-configures colors, tooltip, grid, axes. For standard charts, only write data mapping:
```jsx
<EChart height={{300}} option={{{{ xAxis: {{ type: 'category', data: rows.map(r => r.name) }}, yAxis: {{ type: 'value' }}, series: [{{ type: 'bar', data: rows.map(r => r.val) }}] }}}} />
<EChart height={{300}} option={{{{ tooltip: {{ trigger: 'item' }}, series: [{{ type: 'pie', radius: ['45%','75%'], data: rows.map(r => ({{ value: r.amt, name: r.lbl }})) }}] }}}} />
<EChart height={{300}} option={{{{ xAxis: {{ type: 'category', data: rows.map(r => r.date) }}, yAxis: {{ type: 'value' }}, series: [{{ type: 'line', data: rows.map(r => r.val), areaStyle: {{ opacity: 0.15 }} }}] }}}} />
```
For advanced charts (radar, gauge, treemap, sunburst, funnel, sankey, calendar heatmap, parallel coordinates, graph), pass the full ECharts option — the theme still provides colors and tooltip:
```jsx
<EChart height={{300}} option={{{{ radar: {{ indicator: indicators }}, series: [{{ type: 'radar', data: radarData }}] }}}} />
<EChart height={{250}} option={{{{ series: [{{ type: 'gauge', data: [{{ value: 72 }}], detail: {{ formatter: '{{value}}%' }} }}] }}}} />
<EChart height={{400}} option={{{{ series: [{{ type: 'treemap', data: treeData }}] }}}} />
```

═══════════════════════════════════════════════════════════════════════════════
CONTRAST CONTRACT — NON-NEGOTIABLE (read before writing any styling)
═══════════════════════════════════════════════════════════════════════════════
Pick ONE background mode for the whole dashboard (light OR dark — choose per topic) and theme EVERYTHING to match it. Never dark text on a dark surface or light text on a light surface.
- The `KPICard`/`SectionCard`/`FilterSelect` DEFAULTS (bg-white, text-slate-900) are tuned for a LIGHT background. On a DARK page you MUST pass dark-mode classes to every one of them — e.g. `className="bg-slate-900 border-slate-700 text-slate-100" titleClassName="text-slate-400"`. A dark page with default cards = invisible dark-on-dark text.
- The `<EChart>` 'dash' theme is tuned for LIGHT backgrounds (dark axis labels/legend). On a DARK page, EVERY chart MUST override text to light, or axis labels and legends vanish. Add to each option:
  ```jsx
  textStyle: {{ color: '#e2e8f0' }},
  xAxis: {{ ...xAxis, axisLabel: {{ color: '#cbd5e1' }}, axisLine: {{ lineStyle: {{ color: '#475569' }} }} }},
  yAxis: {{ ...yAxis, axisLabel: {{ color: '#cbd5e1' }}, splitLine: {{ lineStyle: {{ color: '#1e293b' }} }} }},
  legend: {{ textStyle: {{ color: '#cbd5e1' }} }},
  ```
  (On a light page the 'dash' defaults are correct — don't override.)
- Custom `<div>` KPI tiles/tables on a dark page: set explicit light text (`text-slate-100`, muted `text-slate-400`). Don't rely on inherited/default colors.
- After writing the dashboard, mentally check: is any number, label, axis tick, legend, or title the same lightness as the surface behind it? If yes, fix it.

AVAILABLE COMPONENTS (convenience shortcuts — not requirements):
- `<KPICard title="" value={{fmt(n, {{currency:true}})}} subtitle="" color="#3B82F6" className="" titleClassName="" subtitleClassName="" style={{{{}}}} />` — `className` replaces default theme (bg-white, border, text-slate-900). `titleClassName`/`subtitleClassName` replace title/subtitle defaults. `style` for inline overrides. Theme these to match your color story:
  - Dark: `className="bg-slate-900 border-slate-700 text-white" titleClassName="text-slate-400"`
  - Colored: `className="bg-indigo-50 border-indigo-200 text-indigo-900" titleClassName="text-indigo-600"`
- `<SectionCard title="" subtitle="" className="" titleClassName="" subtitleClassName="" style={{{{}}}}>...children...</SectionCard>` — same theming: `className` replaces defaults, `titleClassName`/`subtitleClassName` for text. Theme to match.
- `<FilterSelect label="" options={{arr}} selected={{arr}} onChange={{fn}} searchable={{bool}} className="" style={{{{}}}} />` — multi-select dropdown (portaled). Built-in search at 8+ options. `className` replaces default theme (bg-white border-slate-200 text-slate-900) — pass e.g. `className="bg-slate-900 border-slate-700 text-slate-100"` for dark.
- `<FilterSearch label="" value={{str}} onChange={{e => setFilter(field, e.target.value)}} placeholder="Search..." className="" style={{{{}}}} />` — text search. `className` replaces default theme.
- `<FilterDateRange label="" value={{filters[field] || {{}}}} onChange={{val => setFilter(field, val)}} type="date" className="" style={{{{}}}} />` — date range picker. `className` replaces default theme.
- `fmt(n, opts)` — `{{currency:true}}`, `{{pct:true}}`, auto K/M/B
- `<LoadingSpinner size={{32}} />`

All components are fully themeable via `className`/`titleClassName`/`subtitleClassName`/`style`. Don't leave default white/slate styling when your design calls for something different. If the design needs something these can't express — build custom React + Tailwind.

**INFO POPOVER (required):** Pass `viz={{viz[N]}}` to every `<KPICard>` and `<SectionCard>` you build from a visualization. This renders a small built-in "ⓘ" button that lets users inspect the data behind each component (Data tab with rows, Code tab with the query). Use the index of the visualization the card is derived from (the primary one if it combines several). When a card renders FILTERED rows (you called `filterRows(viz[N].rows)`), ALSO pass `rows={{<those filtered rows>}}` so the popover shows the filtered view that matches the component, not the full dataset. When a card AGGREGATES or derives its value client-side, ALSO pass `calc="<formula>"` describing the math with real column names, e.g. `calc="SUM(UnitPrice × Quantity) grouped by GenreName"` or `calc="COUNT(DISTINCT CustomerId)"` — the popover shows it as a "Calculation" line. If you render a chart with a bare `<EChart>` that is NOT inside a `<SectionCard>`, pass `viz={{viz[N]}}` (and `rows`/`calc` if relevant) to the `<EChart>` itself so it still gets the popover.

**CUSTOM MARKUP — add `data-dash-*` attributes (required):** Whenever you build your OWN containers instead of `<KPICard>`/`<SectionCard>`/`<EChart>` (custom `<div>` KPI tiles, chart wrappers, tables), annotate each item's outer element with `data-dash-viz="N"` (source visualization index) and `data-dash-calc="<formula>"` when the value is derived. A global overlay then renders the same Data/Code/Calc popover on each item. Example: `<div data-dash-viz={{0}} data-dash-calc="SUM(UnitPrice × Quantity)">...custom tile...</div>`. EVERY metric, chart, and table must be reachable via either a prebuilt component's `viz` prop OR a `data-dash-viz` attribute — never leave an item with no way to inspect its data.

DATA ACCESS:

```javascript
const data = useArtifactData(); // Returns null while loading
// data = {{ report: {{id, title}}, visualizations: [...] }}
```

Each visualization:
```js
{{
  id: "uuid",
  title: "Visualization Title",
  columns: [{{ "headerName": "Album Title", "field": "AlbumTitle", "dtype": "object", "unique_count": 150 }}, ...],
  rows: [{{ "AlbumTitle": "Battlestar Galactica", "total_revenue": 35.82 }}, ...],
  view: {{ /* chart config hints */ }},
  dataModel: {{ /* series/axis config */ }}
}}
```

- Use `column.field` to access row values: `row[column.field]`
- Use `column.headerName` for display labels
- Column metadata includes `dtype` (pandas type) and `unique_count` — use these for filter/format decisions
- **Do not hardcode data** — all values should come from `data.visualizations[N].rows`
- **Defensive coding**: Row values and properties can be `null`/`undefined`. Use optional chaining or fallbacks before calling `.includes()`, `.toLowerCase()`, `.startsWith()`, `.split()`, etc. Example: `(row.name || '').includes('x')` or `String(val ?? '').toLowerCase()`. Do not call string methods on a value that could be nullish.

View hints — honor the viz config:
The `view_config` on each visualization describes how the author wants the data rendered. Follow it when generating code.

- `view_config.aggregation` (`"sum" | "avg" | "count" | "min" | "max"`): the raw rows are granular, so aggregate the relevant value column before rendering (especially for `count`, `metric_card`, `pie_chart`, `heatmap`). Use `rows.reduce(...)`. Example for a metric card with aggregation=sum:
  ```js
  const total = useMemo(
    () => viz[0].rows.reduce((s, r) => s + (Number(r.revenue) || 0), 0),
    [viz]
  );
  ```
  For pie/heatmap/bar charts that group by a category, group first and aggregate the value per group rather than using the first matching row.

- `view_config.series_aggregations` (array of `{{key, aggregation}}`): apply the given aggregation per series when building multi-series bar/line/area charts.

- `view_config.default_filters` (array of `{{column, operator, value}}`): the author wants the dashboard to open with these filters already applied. Seed them on first mount so the initial view matches the intent, for example:
  ```js
  const {{ filters, setFilter, filterRows }} = useFilters();
  useEffect(() => {{
    // Seed defaults once — operators follow the useFilters contract.
    {{/* for each entry in view_config.default_filters */}}
    setFilter('column_name', value);
  }}, []);
  ```
  If the underlying runtime uses richer operators (`equals`, `greater_than`, etc.), either call `setFilter` with the operator-aware object it expects, or compute the filtered rows directly via `filterRows(viz[N].rows)` once the filter is seeded. Render the filtered view when defaults are present so the initial numbers match the author's intent.

YOUR VISUALIZATIONS:

{viz_json}

{"(Full sample data included above)" if allow_llm_see_data else "(Data samples hidden for privacy - use column names and row_count to understand the data structure)"}

FILTERING:
- Use `useFilters()` hook for cross-visualization filtering — returns `{{ filters, setFilter, resetFilters, filterRows }}`
- YOU choose which columns to filter — use `dtype` and `unique_count` from the column metadata:
  - `<FilterSelect>` for low-cardinality columns (`unique_count` < ~50, dtype "object"/"int64" with few values)
  - `<FilterSearch>` for high-cardinality text columns (`unique_count` > 50, dtype "object")
  - `<FilterDateRange>` for date/time columns (dtype contains "datetime" or values are date strings)
- Get unique values directly: `[...new Set(viz[N].rows.map(r => r[field]))]`

FILTER FEASIBILITY AUDIT — DO THIS FIRST, BEFORE WRITING CODE:
Before wiring any cross-viz filter, verify it will actually work. A filter that looks wired but silently leaves some vizs untouched is a broken dashboard, not a partial one.

For each dimension you intend to filter by:
1. **Enumerate participating vizs** — which vizs should this filter affect? (Usually: any viz whose topic logically shares the dimension, e.g. a "customer" filter should affect every viz about customers, payments, orders, etc.)
2. **Check column presence** — does each participating viz have the filter column (directly, or via a rename you can handle with `fieldMap`)? Check the `columns` array in YOUR VISUALIZATIONS below.
3. **Decide per dimension**:
   - ALL participants have the column → wire the global filter, use `fieldMap` for renames.
   - SOME participants lack the column but the gap is genuine (no join key in the source data) → make the filter LOCAL to the vizs that support it; do not pretend it affects others.
   - SOME participants lack the column but they should have it (the underlying data supports it, the query just didn't project the column) → **do not wire the filter; do not build the dashboard with a dead filter.** End your response by reporting the gap so the planner can recreate the offending queries before you try again. Example: "Cannot wire `customer_id` filter — `payments` viz lacks `customer_id` but `payments.customer_id` exists in schema. Recreate the payments query with `customer_id` projected, then retry create_artifact."

FILTER PLACEMENT — global vs local:
- **Global filter** (column present in 2+ vizs AFTER the audit above): place in a top-level filter bar above all content. Use one shared filter + `fieldMap` for renames, not duplicates.
- **Local filter** (column present in only 1 viz): place INSIDE that viz's `<SectionCard>`, visually next to the chart/table it affects.
- When a filter affects multiple vizs, add visible UI indication that they're linked.

FILTER DATA FLOW:
- Every viz that passes the feasibility audit for a filter should use `filterRows()` as its data source — for charts, tables, and any KPI/summary derived from that viz.
- KPI cards that summarize filtered data (sum, count, avg) should be computed from filtered rows, not from raw `viz[N].rows`.
- Do not call `filterRows` on a viz that doesn't have the filter column just to "be safe" — silently passing rows through makes the filter look active when it isn't. Audit first, wire second.

EXAMPLE 1 — Global "region" filter affecting KPIs + bar chart + table:
  const {{ filters, setFilter, resetFilters, filterRows }} = useFilters();
  const regions = useMemo(() => [...new Set(vizSales.rows.map(r => r.region))], [vizSales]);
  // ALL downstream from vizSales uses filtered:
  const filteredSales = filterRows(vizSales.rows);
  const totalRevenue = useMemo(() => filteredSales.reduce((s, r) => s + r.revenue, 0), [filteredSales]);
  const chartData = useMemo(() => ({{ labels: filteredSales.map(r => r.month), values: filteredSales.map(r => r.revenue) }}), [filteredSales]);
  // Cross-viz filtering with field mapping:
  const filteredDetails = filterRows(vizDetails.rows, {{ region: 'RegionName' }});
  // Layout: <FilterSelect> in top bar, KPIs below, charts below that

EXAMPLE 2 — Local filter inside a SectionCard:
  const {{ filters, setFilter, filterRows }} = useFilters();
  const filtered = filterRows(vizProducts.rows);
  // Layout: <SectionCard title="Products"><FilterSelect .../><EChart ... /></SectionCard>

- Include a Reset button when any filters are active (`Object.keys(filters).length > 0`)
- After filtering, if a visualization has zero matching rows, display "No data matches current filters"

═══════════════════════════════════════════════════════════════════════════════
DESIGN GUIDANCE (use when the user hasn't specified styling)
═══════════════════════════════════════════════════════════════════════════════

If the user specified a theme/style/colors above, follow that — skip this section.
Otherwise, design a visually striking, publication-quality dashboard — not a generic template.

COLOR & IDENTITY:
- Pick a cohesive color story that fits the data topic. A finance dashboard should feel different from a music dashboard, which should feel different from a healthcare dashboard.
- Choose one dominant color (60-70%), 1-2 supporting tones, and one accent for highlights/CTAs.
- Do NOT default to generic blue. Blue is fine if it fits the topic — but earn it, don't default to it.
- Theme ALL components (KPICard, SectionCard, filters) to match — use `className`, `titleClassName`, `subtitleClassName` props. Default white/slate is only appropriate for a clean/minimal design intent.

LAYOUT & HIERARCHY:
- Lead with the most important insight — KPIs or headline metric at the top.
- Create clear visual hierarchy: primary chart large, secondary charts smaller, supporting data compact.
- Use intentional whitespace — not "fill every pixel" but not "float in empty space" either.
- Vary card sizes and chart heights to create rhythm. A grid of same-sized boxes is boring.

TYPOGRAPHY & POLISH:
- Clean, modern typography. Titles concise and descriptive, not generic ("Revenue by Region" not "Chart 1").
- Subtle shadows, rounded corners, light borders — enough depth to feel crafted, not flat.
- Light mode default. Dark mode only if the topic or user suggests it.

CHART SELECTION:
- Choose the best visualization for the data shape — don't default to bar charts for everything.
- Standard charts (bar, line, pie, area) for simple relationships. Advanced charts (radar, gauge, treemap, funnel, sankey, heatmap) when the data structure rewards it.
- Show data from different angles without redundancy. Each chart should reveal something the others don't.

The goal: it should look like a designer built it for this specific dataset, not like a template was filled in.

═══════════════════════════════════════════════════════════════════════════════
OUTPUT FORMAT
═══════════════════════════════════════════════════════════════════════════════

```
<script type="text/babel">
function App() {{
  const data = useArtifactData();
  if (!data) return <div className="flex items-center justify-center h-screen text-gray-400"><LoadingSpinner size={{32}} /></div>;
  const viz = data.visualizations;
  // ... concise dashboard code
}}
ReactDOM.createRoot(document.getElementById('root')).render(<App />);
</script>
```

Structure: all code should be inside `function App() {{ ... }}` with `ReactDOM.createRoot(document.getElementById('root')).render(<App />);` at the end. Do not put return statements outside a function.

Rules: `<script type="text/babel">` wrapper. `useArtifactData()` for data. `<EChart option={{...}} />` for charts. Pass `viz={{viz[N]}}` to every KPICard/SectionCard so the built-in info popover shows the data behind it. Responsive. Handle zero rows. No hardcoded data. No UUIDs/branding/emoji. Guard nullish values before string methods (use `(val || '')` or `String(val ?? '')`).

**Code size:** Write compact code — no unnecessary variables, comments, or verbose JSX. Omit default props. Don't repeat theme styling the 'dash' theme already provides. Prefer inline expressions over separate variables when used once. For simple dashboards target under 8K characters. For detailed/specific user requests, use as much space as needed to faithfully implement their design — fidelity to the user's request is more important than brevity.

Now create the dashboard:"""

    def _build_prompt(
        self,
        user_prompt: str,
        title: str | None,
        mode: str,
        viz_profiles: List[Dict[str, Any]],
        instructions_context: str,
        report_title: str | None,
        allow_llm_see_data: bool,
        messages_context: str = "",
        image_count: int = 0,
        organization_settings: Any = None,
    ) -> str:
        """Build the prompt for generating artifact code. Dispatches to mode-specific builders."""
        if mode == "slides":
            return self._build_slides_prompt(
                user_prompt=user_prompt,
                title=title,
                viz_profiles=viz_profiles,
                instructions_context=instructions_context,
                report_title=report_title,
                allow_llm_see_data=allow_llm_see_data,
                messages_context=messages_context,
                image_count=image_count,
                organization_settings=organization_settings,
            )
        return self._build_page_prompt(
            user_prompt=user_prompt,
            title=title,
            viz_profiles=viz_profiles,
            instructions_context=instructions_context,
            report_title=report_title,
            allow_llm_see_data=allow_llm_see_data,
            messages_context=messages_context,
            image_count=image_count,
            organization_settings=organization_settings,
        )

    def _extract_code(self, response: str, mode: str = "page") -> str:
        """Extract the code from the LLM response.

        For 'page' mode: Extract React code from <script type="text/babel"> tags
        For 'slides' mode: Extract python-pptx code from python code blocks
        """
        if mode == "slides":
            return self._extract_slides_python(response)

        # Dashboard mode - extract React code from script tags
        start_marker = "<script type=\"text/babel\">"
        end_marker = "</script>"

        start_idx = response.find(start_marker)
        if start_idx == -1:
            # Try alternative markers
            start_marker = "<script type='text/babel'>"
            start_idx = response.find(start_marker)

        if start_idx != -1:
            end_idx = response.find(end_marker, start_idx)
            if end_idx != -1:
                code = response[start_idx:end_idx + len(end_marker)]
                return self._sanitize_code(self._ensure_app_wrapper(code))

        # If no script tags found, wrap the response
        code = response.strip()
        if not code.startswith("<script"):
            code = f'<script type="text/babel">\n{code}\n</script>'

        return self._sanitize_code(self._ensure_app_wrapper(code))

    @staticmethod
    def _sanitize_code(code: str) -> str:
        """Fix common LLM code generation artifacts deterministically."""
        import re

        # Strip ES-module syntax. The sandbox provides React/ReactDOM/echarts/the
        # EChart component/useArtifactData() as GLOBALS and runs the artifact via
        # babel-standalone in a <script type="text/babel"> (NO module support).
        # Any `import`/`export` the model emits despite the no-import contract
        # makes the browser throw "Cannot use import statement outside a module",
        # which fails the whole dashboard render. Since the libs are already
        # global these statements are always redundant — remove them defensively
        # so a non-compliant model can't break the render.
        #   import X from '...';  / import {a,b} from '...';  / import * as ns from '...';  (incl. multiline)
        code = re.sub(r'^[ \t]*import\s[\s\S]*?from\s*[\'"][^\'"]+[\'"]\s*;?[ \t]*$', '', code, flags=re.M)
        #   bare side-effect import:  import '...';
        code = re.sub(r'^[ \t]*import\s+[\'"][^\'"]+[\'"]\s*;?[ \t]*$', '', code, flags=re.M)
        #   export default <decl>  ->  <decl>   (keep the declaration, drop the keyword)
        code = re.sub(r'^([ \t]*)export\s+default\s+', r'\1', code, flags=re.M)
        #   export const/let/var/function/class  ->  drop the `export ` prefix
        code = re.sub(r'^([ \t]*)export\s+(?=(?:const|let|var|function|class|async)\b)', r'\1', code, flags=re.M)
        #   standalone  export { ... };  -> remove the line entirely
        code = re.sub(r'^[ \t]*export\s*\{[^}]*\}\s*;?[ \t]*$', '', code, flags=re.M)

        # Fix double-brace pattern: function App() {\n{ ... }\n}
        # The LLM sometimes wraps the function body in an extra block scope.
        # Match: function App() {\n{ at the start, and }\n} at the end (before render call)
        code = re.sub(
            r'(function\s+\w+\s*\([^)]*\)\s*\{)\s*\n\s*\{',
            r'\1',
            code,
        )
        # Remove the matching trailing extra }
        # Look for }\n}\n before ReactDOM.createRoot
        code = re.sub(
            r'\}\s*\n\s*\}\s*\n(\s*ReactDOM\.createRoot)',
            r'}\n\1',
            code,
        )

        return code

    @staticmethod
    def _ensure_app_wrapper(code: str) -> str:
        """Ensure code has a proper App component wrapper.

        LLM sometimes outputs bare return statements outside a function.
        Detect and fix by wrapping the inner code in function App() + ReactDOM.createRoot.
        """
        import re

        # Check if code already has an App function/component
        if re.search(r'function\s+App\s*\(', code) or re.search(r'(?:const|let|var)\s+App\s*=', code):
            return code

        # Extract inner code between script tags
        inner_match = re.search(
            r'<script\s+type=["\']text/babel["\']>\s*([\s\S]*?)\s*</script>',
            code
        )
        if not inner_match:
            return code

        inner = inner_match.group(1).strip()

        # Strip any existing broken ReactDOM.createRoot/render calls
        inner = re.sub(r'ReactDOM\.createRoot\(.*?\)\.render\(.*?\);?\s*$', '', inner, flags=re.DOTALL).strip()

        logger.warning("_ensure_app_wrapper: LLM output missing function App() wrapper — auto-wrapping")

        wrapped = (
            '<script type="text/babel">\n'
            'function App() {\n'
            f'{inner}\n'
            '}\n'
            "ReactDOM.createRoot(document.getElementById('root')).render(<App />);\n"
            '</script>'
        )
        return wrapped

    def _extract_slides_python(self, response: str) -> str:
        """Extract python-pptx code for slides mode."""
        import re

        # Try to find Python code block
        python_match = re.search(r'```python\s*([\s\S]*?)```', response)
        if python_match:
            return python_match.group(1).strip()

        # Try generic code block
        code_match = re.search(r'```\s*([\s\S]*?)```', response)
        if code_match:
            return code_match.group(1).strip()

        # Look for function definition as start marker
        func_start = response.find('def generate_slides')
        if func_start != -1:
            # Find the prs.save() call at the end
            save_end = response.rfind('prs.save(')
            if save_end != -1:
                # Include the full save line
                end_idx = response.find(')', save_end)
                if end_idx != -1:
                    return response[func_start:end_idx + 1].strip()
            return response[func_start:].strip()

        # Fallback: return the response as-is
        return response.strip()

