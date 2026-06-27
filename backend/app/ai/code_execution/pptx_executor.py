"""
Secure executor for python-pptx code generation.

Allows LLM-generated python-pptx code to run in a sandboxed environment
with security validation reused from code_execution.py.
"""

import io
import ast
import tempfile
import subprocess
from pathlib import Path
from contextlib import redirect_stdout
from typing import Dict, Any, List, Tuple, Optional

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, ChartData

from app.ai.code_execution.code_execution import (
    CodeSecurityError,
    UnsafePythonError,
    FORBIDDEN_MODULES,
    FORBIDDEN_BUILTINS,
    FORBIDDEN_ATTRIBUTES,
)


# =============================================================================
# PPTX-Specific Security Validation
# =============================================================================

# Modules allowed for PPTX generation (extend the forbidden list with exceptions)
PPTX_ALLOWED_MODULES = frozenset({
    'pptx',
})

# Safe read-only introspection builtins the pptx codegen legitimately needs.
# python-pptx styling is full of optional attributes (category_axis, value_axis,
# has_data_labels…) so the slides prompt mandates a `style_chart_text` helper
# built on getattr/hasattr. These are read-only attribute access — no exec/eval/
# import power — so we whitelist them out of FORBIDDEN_BUILTINS for this sandbox
# ONLY. setattr is deliberately NOT whitelisted (it could mutate dunders).
PPTX_ALLOWED_BUILTINS = frozenset({
    'getattr',
    'hasattr',
})


class PptxSecurityVisitor(ast.NodeVisitor):
    """AST visitor that checks for dangerous code patterns, allowing pptx imports."""

    def __init__(self):
        self.errors: List[str] = []

    def visit_Import(self, node: ast.Import):
        for alias in node.names:
            module_name = alias.name.split('.')[0]
            # Allow pptx imports, block everything else that's forbidden
            if module_name not in PPTX_ALLOWED_MODULES and module_name in FORBIDDEN_MODULES:
                self.errors.append(f"Forbidden import: '{alias.name}'")
        self.generic_visit(node)

    def visit_ImportFrom(self, node: ast.ImportFrom):
        if node.module:
            module_name = node.module.split('.')[0]
            # Allow pptx imports, block everything else that's forbidden
            if module_name not in PPTX_ALLOWED_MODULES and module_name in FORBIDDEN_MODULES:
                self.errors.append(f"Forbidden import: 'from {node.module}'")
        self.generic_visit(node)

    def visit_Call(self, node: ast.Call):
        # Check for forbidden built-in calls like eval(), exec(), open()
        if isinstance(node.func, ast.Name):
            if node.func.id in FORBIDDEN_BUILTINS and node.func.id not in PPTX_ALLOWED_BUILTINS:
                self.errors.append(f"Forbidden function call: '{node.func.id}()'")

        # Check for __import__('os') style calls
        if isinstance(node.func, ast.Name) and node.func.id == '__import__':
            self.errors.append("Forbidden function call: '__import__()'")

        self.generic_visit(node)

    def visit_Attribute(self, node: ast.Attribute):
        # Check for direct access to forbidden attributes like obj.__class__
        if node.attr in FORBIDDEN_ATTRIBUTES:
            self.errors.append(f"Forbidden attribute access: '{node.attr}'")
        self.generic_visit(node)


def validate_pptx_code(code: str) -> None:
    """
    Validate Python code for security issues, allowing pptx imports.

    Raises:
        UnsafePythonError: If the code contains dangerous constructs.
    """
    try:
        tree = ast.parse(code)
    except SyntaxError:
        # Let syntax errors pass through - they'll fail at exec() time
        return

    visitor = PptxSecurityVisitor()
    visitor.visit(tree)

    if visitor.errors:
        raise UnsafePythonError(
            f"Code contains forbidden constructs: {'; '.join(visitor.errors)}"
        )


# =============================================================================
# PPTX Code Executor
# =============================================================================

class PptxCodeExecutor:
    """
    Secure executor for python-pptx code generation.

    Reuses security patterns from StreamingCodeExecutor but with a namespace
    tailored for PPTX generation.
    """

    def __init__(self, logger=None):
        self.logger = logger

    def execute_pptx_code(
        self,
        *,
        code: str,
        visualizations: List[Dict[str, Any]],
        report: Dict[str, Any],
        output_path: Path,
    ) -> Tuple[Path, str]:
        """
        Execute python-pptx code and save the resulting presentation.

        Args:
            code: The python-pptx code to execute
            visualizations: List of visualization dicts with rows/columns
            report: Report info dict with id, title, theme
            output_path: Path where the PPTX file should be saved

        Returns:
            Tuple of (output_path, stdout_log)

        Raises:
            UnsafePythonError: If code contains forbidden imports, calls, or attributes
        """
        # Security: Validate code before execution
        validate_pptx_code(code)

        output_log = ""

        # Build the namespace with pptx utilities and data
        local_namespace = {
            # python-pptx classes
            'Presentation': Presentation,
            'Inches': Inches,
            'Pt': Pt,
            'Emu': Emu,
            'RGBColor': RGBColor,
            'PP_ALIGN': PP_ALIGN,
            'MSO_ANCHOR': MSO_ANCHOR,
            'MSO_SHAPE': MSO_SHAPE,
            'XL_CHART_TYPE': XL_CHART_TYPE,
            'XL_LEGEND_POSITION': XL_LEGEND_POSITION,
            'CategoryChartData': CategoryChartData,
            'ChartData': ChartData,

            # Data access
            'visualizations': visualizations,
            'report': report,

            # Output target (set by executor, not user code)
            '_pptx_output_path': str(output_path),
        }

        if self.logger:
            self.logger.debug(f"Executing PPTX code:\n{code}")

        with io.StringIO() as stdout_capture:
            with redirect_stdout(stdout_capture):
                exec(code, local_namespace)
            output_log = stdout_capture.getvalue()

        # Verify the file was created. Models frequently wrap the whole deck in a
        # `def generate_slides(...)` and either never call it, or call prs.save()
        # with a hardcoded filename instead of _pptx_output_path. Both leave no
        # file at output_path. Rather than hard-fail (which dumps raw code to the
        # UI), attempt a deterministic rescue before giving up.
        if not output_path.exists():
            rescued = self._rescue_save(
                local_namespace, output_path, visualizations, report
            )
            if rescued and self.logger:
                self.logger.warning(
                    "PPTX code did not save to _pptx_output_path; rescued the "
                    "Presentation object and saved it for the user."
                )

        if not output_path.exists():
            raise RuntimeError(
                f"PPTX code executed but no file was created at {output_path}. "
                "Ensure the code calls prs.save(_pptx_output_path)"
            )

        return output_path, output_log

    def _rescue_save(
        self,
        namespace: Dict[str, Any],
        output_path: Path,
        visualizations: List[Dict[str, Any]],
        report: Dict[str, Any],
    ) -> bool:
        """Best-effort recovery when generated code didn't save to output_path.

        Two common failure modes:
          1. The deck is built inside a function (generate_slides/build/main/...)
             that the code never calls — so nothing ran. Invoke it.
          2. A Presentation object exists in the namespace but was never saved,
             or was saved to the wrong path. Save it to output_path.

        Returns True if a file now exists at output_path.
        """
        def _is_presentation(obj) -> bool:
            # python-pptx Presentation() returns an instance with .save + .slides
            return hasattr(obj, "save") and hasattr(obj, "slides") and callable(
                getattr(obj, "save", None)
            )

        # Mode 1: call a conventional entrypoint that was defined but never invoked.
        for fn_name in (
            "generate_slides", "build_slides", "build_presentation",
            "create_presentation", "build", "main", "make_deck",
        ):
            fn = namespace.get(fn_name)
            if not callable(fn) or output_path.exists():
                continue
            # Try the most common signatures, widest first.
            for call_args in (
                (visualizations, report),
                (visualizations,),
                (),
            ):
                try:
                    ret = fn(*call_args)
                except TypeError:
                    continue  # signature mismatch — try the next arity
                except Exception:
                    break  # ran but blew up internally — don't keep retrying
                # If it handed back a Presentation, persist it.
                if _is_presentation(ret):
                    try:
                        ret.save(str(output_path))
                    except Exception:
                        pass
                break

        # Mode 2: scan the namespace for any Presentation object and save it.
        if not output_path.exists():
            for val in namespace.values():
                if _is_presentation(val):
                    try:
                        val.save(str(output_path))
                        break
                    except Exception:
                        continue

        return output_path.exists()

    def execute_with_retries(
        self,
        *,
        code: str,
        visualizations: List[Dict[str, Any]],
        report: Dict[str, Any],
        output_path: Path,
        fix_code_fn: Optional[callable] = None,
        max_retries: int = 2,
    ) -> Tuple[Path, str, List[Tuple[str, str]]]:
        """
        Execute PPTX code with retry logic.

        Args:
            code: Initial python-pptx code
            visualizations: Visualization data
            report: Report info
            output_path: Output path for PPTX
            fix_code_fn: Optional async function to fix code on error
            max_retries: Maximum number of retry attempts

        Returns:
            Tuple of (output_path, stdout_log, code_and_error_messages)
        """
        code_and_error_messages: List[Tuple[str, str]] = []
        current_code = code

        for attempt in range(max_retries):
            try:
                result_path, output_log = self.execute_pptx_code(
                    code=current_code,
                    visualizations=visualizations,
                    report=report,
                    output_path=output_path,
                )
                return result_path, output_log, code_and_error_messages
            except Exception as e:
                error_msg = str(e)
                code_and_error_messages.append((current_code, error_msg))

                if self.logger:
                    self.logger.warning(
                        f"PPTX execution attempt {attempt + 1} failed: {error_msg}"
                    )

                # If we have a fix function and more retries, try to fix
                if fix_code_fn and attempt < max_retries - 1:
                    try:
                        current_code = fix_code_fn(current_code, error_msg)
                    except Exception as fix_error:
                        if self.logger:
                            self.logger.error(f"Code fix failed: {fix_error}")
                        raise e
                else:
                    raise e

        # Should not reach here, but just in case
        raise RuntimeError("Max retries exceeded for PPTX code execution")


# =============================================================================
# PPTX to Image Preview Conversion
# =============================================================================

class PptxPreviewService:
    """
    Service for generating preview images from PPTX files.

    Uses LibreOffice (headless) to convert PPTX to PDF,
    then pdf2image to convert PDF pages to PNG images.
    """

    def __init__(self, preview_dir: Optional[Path] = None, logger=None):
        self.logger = logger
        if preview_dir:
            self.preview_dir = preview_dir
        else:
            # Default to uploads/pptx_previews relative to backend root
            backend_root = Path(__file__).parent.parent.parent.parent
            self.preview_dir = backend_root / "uploads" / "pptx_previews"

        # Ensure preview directory exists
        self.preview_dir.mkdir(parents=True, exist_ok=True)

    def generate_previews(
        self,
        pptx_path: Path,
        artifact_id: str,
        dpi: int = 150,
    ) -> List[str]:
        """
        Convert PPTX to PNG preview images.

        Args:
            pptx_path: Path to the PPTX file
            artifact_id: Artifact ID for organizing previews
            dpi: Resolution for preview images (default 150)

        Returns:
            List of relative paths to preview images (e.g., ["pptx_previews/{id}/slide-1.png", ...])
        """
        from pdf2image import convert_from_path
        from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError

        # Create artifact-specific preview directory
        artifact_preview_dir = self.preview_dir / artifact_id
        artifact_preview_dir.mkdir(parents=True, exist_ok=True)

        # Step 1: Convert PPTX to PDF using LibreOffice
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)

            # LibreOffice convert to PDF
            try:
                result = subprocess.run(
                    [
                        'soffice',
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', str(tmp_path),
                        str(pptx_path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=60,
                )
                if result.returncode != 0:
                    raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
            except FileNotFoundError:
                raise RuntimeError(
                    "LibreOffice not found. Install with: apt-get install libreoffice-impress"
                )

            # Find the generated PDF
            pdf_files = list(tmp_path.glob("*.pdf"))
            if not pdf_files:
                raise RuntimeError("LibreOffice did not produce a PDF file")
            pdf_path = pdf_files[0]

            # Step 2: Convert PDF pages to PNG using pdf2image
            try:
                images = convert_from_path(pdf_path, dpi=dpi)
            except PDFInfoNotInstalledError:
                raise RuntimeError(
                    "poppler not found. Install with: apt-get install poppler-utils"
                )
            except PDFPageCountError as e:
                raise RuntimeError(f"Failed to read PDF: {e}")

            # Save images
            for i, image in enumerate(images):
                image_path = artifact_preview_dir / f"slide-{i + 1:02d}.png"
                image.save(str(image_path), "PNG")

        # Collect generated image paths
        preview_images = sorted(artifact_preview_dir.glob("slide-*.png"))

        # Return relative paths from uploads directory
        relative_paths = [
            f"pptx_previews/{artifact_id}/{img.name}" for img in preview_images
        ]

        if self.logger:
            self.logger.info(f"Generated {len(relative_paths)} preview images for artifact {artifact_id}")

        return relative_paths

    def get_preview_paths(self, artifact_id: str) -> List[str]:
        """Get existing preview image paths for an artifact."""
        artifact_preview_dir = self.preview_dir / artifact_id
        if not artifact_preview_dir.exists():
            return []

        preview_images = sorted(artifact_preview_dir.glob("slide-*.png"))
        return [
            f"pptx_previews/{artifact_id}/{img.name}" for img in preview_images
        ]

    def cleanup_previews(self, artifact_id: str) -> None:
        """Remove all preview images for an artifact."""
        artifact_preview_dir = self.preview_dir / artifact_id
        if artifact_preview_dir.exists():
            import shutil
            shutil.rmtree(artifact_preview_dir)
